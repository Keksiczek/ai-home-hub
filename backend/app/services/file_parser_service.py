"""File parser service – extract text from various file formats with OCR support."""

import logging
import re
import zipfile
from pathlib import Path
from typing import Any, Dict

logger = logging.getLogger(__name__)


class FileParserService:
    """Parse text content from various file formats including OCR for images."""

    SUPPORTED_EXTENSIONS = {
        ".pdf",
        ".docx",
        ".xlsx",
        ".pptx",
        ".txt",
        ".md",
        ".jpg",
        ".jpeg",
        ".png",
        ".gif",
        ".bmp",
        ".mp3",
        ".wav",
        ".m4a",
        ".ogg",
        ".mp4",
        ".webm",
        ".mov",
        ".epub",
        ".html",
        ".htm",
        ".zip",
    }

    # Map extensions to media_type for unified metadata
    MEDIA_TYPES: Dict[str, str] = {
        ".pdf": "text",
        ".docx": "office",
        ".xlsx": "office",
        ".pptx": "office",
        ".txt": "text",
        ".md": "text",
        ".jpg": "image",
        ".jpeg": "image",
        ".png": "image",
        ".gif": "image",
        ".bmp": "image",
        ".mp3": "audio",
        ".wav": "audio",
        ".m4a": "audio",
        ".ogg": "audio",
        ".mp4": "video",
        ".webm": "video",
        ".mov": "video",
        ".epub": "text",
        ".html": "text",
        ".htm": "text",
        ".zip": "archive",
    }

    # Map extensions to MIME types
    MIME_TYPES: Dict[str, str] = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
        ".mp3": "audio/mpeg",
        ".wav": "audio/wav",
        ".m4a": "audio/mp4",
        ".ogg": "audio/ogg",
        ".mp4": "video/mp4",
        ".webm": "video/webm",
        ".mov": "video/quicktime",
        ".epub": "application/epub+zip",
        ".html": "text/html",
        ".htm": "text/html",
        ".zip": "application/zip",
    }

    def parse_file(self, file_path: Path) -> Dict[str, Any]:
        """
        Parse file and extract text content.

        Returns:
            {
                'text': str,
                'metadata': dict,
                'page_count': int,
                'media_type': str,     # text|image|audio|video|office|archive
                'duration_seconds': float | None,  # for audio/video
                'error': str,          # only present on failure
            }
        """
        if not file_path.exists():
            return {"error": f"File not found: {file_path}"}

        ext = file_path.suffix.lower()

        try:
            if ext == ".pdf":
                return self._parse_pdf(file_path)
            elif ext == ".docx":
                return self._parse_docx(file_path)
            elif ext == ".xlsx":
                return self._parse_xlsx(file_path)
            elif ext == ".pptx":
                return self._parse_pptx(file_path)
            elif ext in {".txt", ".md"}:
                return self._parse_text(file_path)
            elif ext in {".jpg", ".jpeg", ".png", ".gif", ".bmp"}:
                return self._parse_image(file_path)
            elif ext in {".mp3", ".wav", ".m4a", ".ogg"}:
                return self._parse_audio(file_path)
            elif ext in {".mp4", ".webm", ".mov"}:
                return self._parse_video(file_path)
            elif ext == ".epub":
                return self._parse_epub(file_path)
            elif ext in {".html", ".htm"}:
                return self._parse_html(file_path)
            elif ext == ".zip":
                return self._parse_zip(file_path)
            else:
                return {"error": f"Unsupported file type: {ext}"}

        except Exception as exc:
            logger.error("Failed to parse %s: %s", file_path, exc)
            return {"error": str(exc)}

    # ── PDF ──────────────────────────────────────────────────────

    def _parse_pdf(self, file_path: Path) -> Dict[str, Any]:
        import PyPDF2

        text_parts = []
        metadata: Dict[str, Any] = {}

        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            page_count = len(reader.pages)

            if reader.metadata:
                metadata = {
                    "author": reader.metadata.get("/Author", ""),
                    "title": reader.metadata.get("/Title", ""),
                    "created": reader.metadata.get("/CreationDate", ""),
                }

            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

        return {
            "text": "\n\n".join(text_parts),
            "metadata": metadata,
            "page_count": page_count,
        }

    # ── DOCX ─────────────────────────────────────────────────────

    def _parse_docx(self, file_path: Path) -> Dict[str, Any]:
        import docx

        doc = docx.Document(file_path)

        text_parts = [para.text for para in doc.paragraphs if para.text.strip()]

        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)

        metadata: Dict[str, Any] = {}
        if doc.core_properties:
            metadata = {
                "author": doc.core_properties.author or "",
                "title": doc.core_properties.title or "",
                "created": (
                    str(doc.core_properties.created)
                    if doc.core_properties.created
                    else ""
                ),
            }

        return {
            "text": "\n\n".join(text_parts),
            "metadata": metadata,
            "page_count": len(doc.sections),
        }

    # ── XLSX ─────────────────────────────────────────────────────

    def _parse_xlsx(self, file_path: Path) -> Dict[str, Any]:
        import openpyxl

        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"=== Sheet: {sheet_name} ===")

            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(
                    str(cell) if cell is not None else "" for cell in row
                )
                if row_text.strip():
                    text_parts.append(row_text)

        metadata = {
            "sheets": wb.sheetnames,
            "sheet_count": len(wb.sheetnames),
        }

        return {
            "text": "\n".join(text_parts),
            "metadata": metadata,
            "page_count": len(wb.sheetnames),
        }

    # ── PPTX ─────────────────────────────────────────────────────

    def _parse_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from PowerPoint presentation slides."""
        from pptx import Presentation

        prs = Presentation(str(file_path))
        text_parts = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            slide_texts.append(para_text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        if row_text.strip():
                            slide_texts.append(row_text)
            if slide_texts:
                text_parts.append(
                    f"=== Slide {slide_idx} ===\n" + "\n".join(slide_texts)
                )

        slide_count = len(prs.slides)
        metadata: Dict[str, Any] = {"slide_count": slide_count}

        return {
            "text": "\n\n".join(text_parts),
            "metadata": metadata,
            "page_count": slide_count,
        }

    # ── TXT / MD ─────────────────────────────────────────────────

    def _parse_text(self, file_path: Path) -> Dict[str, Any]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

        return {
            "text": text,
            "metadata": {},
            "page_count": 1,
        }

    # ── Images with OCR ──────────────────────────────────────────

    def _parse_image(self, file_path: Path) -> Dict[str, Any]:
        """Parse image file with OCR text extraction using Tesseract."""
        from PIL import Image

        img = Image.open(file_path)
        metadata: Dict[str, Any] = {
            "format": str(img.format or ""),
            "width": img.size[0],
            "height": img.size[1],
            "mode": img.mode,
            "type": "image",
        }

        # Try OCR with pytesseract
        ocr_text = self._ocr_tesseract(file_path)

        if ocr_text and ocr_text.strip():
            metadata["ocr_method"] = "tesseract"
            metadata["ocr_text_length"] = len(ocr_text)
            return {
                "text": ocr_text.strip(),
                "metadata": metadata,
                "page_count": 1,
            }

        # Fallback: return image description without OCR text
        metadata["ocr_method"] = "none"
        return {
            "text": f"[Image: {file_path.name}, {img.size[0]}x{img.size[1]}, no text extracted]",
            "metadata": metadata,
            "page_count": 1,
        }

    def _ocr_tesseract(self, file_path: Path) -> str:
        """Extract text from image using Tesseract OCR. Returns empty string on failure."""
        try:
            import pytesseract
            from PIL import Image

            img = Image.open(file_path)
            text = pytesseract.image_to_string(img, lang="eng+ces", timeout=30)
            return text
        except ImportError:
            logger.warning(
                "pytesseract not installed – OCR skipped for %s", file_path.name
            )
            return ""
        except Exception as exc:
            logger.warning("OCR failed for %s: %s", file_path.name, exc)
            return ""

    def parse_image_ocr(self, file_path: str) -> Dict[str, Any]:
        """Public OCR method for external use. Returns {text, metadata, ocr_method}."""
        path = Path(file_path)
        if not path.exists():
            return {
                "text": "",
                "metadata": {},
                "ocr_method": "error",
                "error": "File not found",
            }
        return self._parse_image(path)

    # ── Audio (speech-to-text via faster-whisper) ─────────────────

    def _parse_audio(self, file_path: Path) -> Dict[str, Any]:
        """Transcribe audio using faster-whisper. Falls back to metadata-only."""
        duration_seconds: float | None = None
        metadata: Dict[str, Any] = {"media_type": "audio"}

        # Try to get duration via ffprobe
        duration_seconds = self._get_media_duration(file_path)
        if duration_seconds is not None:
            metadata["duration_seconds"] = duration_seconds

        try:
            from faster_whisper import WhisperModel

            model = WhisperModel("base", device="cpu", compute_type="int8")
            segments, info = model.transcribe(str(file_path), beam_size=5)
            text_parts = [segment.text for segment in segments]
            text = " ".join(text_parts).strip()

            metadata["language"] = getattr(info, "language", None)
            metadata["duration_seconds"] = getattr(info, "duration", duration_seconds)

            if not text:
                return {
                    "text": f"[Audio: {file_path.name}, no speech detected]",
                    "metadata": metadata,
                    "page_count": 1,
                    "duration_seconds": metadata.get("duration_seconds"),
                }

            return {
                "text": text,
                "metadata": metadata,
                "page_count": 1,
                "duration_seconds": metadata.get("duration_seconds"),
            }
        except ImportError:
            logger.warning(
                "faster-whisper not installed – audio transcription skipped for %s",
                file_path.name,
            )
            return {
                "text": f"[Audio: {file_path.name}, transcription unavailable (faster-whisper not installed)]",
                "metadata": metadata,
                "page_count": 1,
                "duration_seconds": duration_seconds,
            }
        except Exception as exc:
            logger.warning("Audio transcription failed for %s: %s", file_path.name, exc)
            return {
                "text": f"[Audio: {file_path.name}, transcription failed: {exc}]",
                "metadata": metadata,
                "page_count": 1,
                "duration_seconds": duration_seconds,
            }

    # ── Video (audio track transcription) ─────────────────────────

    def _parse_video(self, file_path: Path) -> Dict[str, Any]:
        """Extract audio track from video and transcribe via faster-whisper."""
        import tempfile

        metadata: Dict[str, Any] = {"media_type": "video"}
        duration_seconds = self._get_media_duration(file_path)
        if duration_seconds is not None:
            metadata["duration_seconds"] = duration_seconds

        # Extract audio track to temp WAV
        try:
            import subprocess

            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
                tmp_path = Path(tmp.name)

            result = subprocess.run(
                [
                    "ffmpeg",
                    "-i",
                    str(file_path),
                    "-vn",
                    "-acodec",
                    "pcm_s16le",
                    "-ar",
                    "16000",
                    "-ac",
                    "1",
                    str(tmp_path),
                    "-y",
                ],
                capture_output=True,
                timeout=120,
            )
            if result.returncode != 0:
                logger.warning("ffmpeg audio extraction failed for %s", file_path.name)
                return {
                    "text": f"[Video: {file_path.name}, audio extraction failed]",
                    "metadata": metadata,
                    "page_count": 1,
                    "duration_seconds": duration_seconds,
                }

            # Transcribe the extracted audio
            audio_result = self._parse_audio(tmp_path)
            audio_result["metadata"].update(metadata)
            audio_result["metadata"]["media_type"] = "video"
            return audio_result
        except FileNotFoundError:
            logger.warning(
                "ffmpeg not found – video transcription skipped for %s", file_path.name
            )
            return {
                "text": f"[Video: {file_path.name}, transcription unavailable (ffmpeg not found)]",
                "metadata": metadata,
                "page_count": 1,
                "duration_seconds": duration_seconds,
            }
        except Exception as exc:
            logger.warning("Video parse failed for %s: %s", file_path.name, exc)
            return {
                "text": f"[Video: {file_path.name}, parse failed: {exc}]",
                "metadata": metadata,
                "page_count": 1,
                "duration_seconds": duration_seconds,
            }
        finally:
            try:
                tmp_path.unlink(missing_ok=True)
            except Exception:
                pass

    # ── ePub ──────────────────────────────────────────────────────

    def _parse_epub(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from ePub ebook."""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup

            book = epub.read_epub(str(file_path), options={"ignore_ncx": True})
            text_parts = []

            metadata: Dict[str, Any] = {}
            title = book.get_metadata("DC", "title")
            if title:
                metadata["title"] = title[0][0]
            creator = book.get_metadata("DC", "creator")
            if creator:
                metadata["author"] = creator[0][0]
            lang = book.get_metadata("DC", "language")
            if lang:
                metadata["language"] = lang[0][0]

            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    text = soup.get_text(separator="\n", strip=True)
                    if text.strip():
                        text_parts.append(text)

            chapter_count = len(text_parts)
            metadata["chapter_count"] = chapter_count

            return {
                "text": "\n\n".join(text_parts),
                "metadata": metadata,
                "page_count": chapter_count or 1,
            }
        except ImportError:
            logger.warning(
                "ebooklib not installed – ePub parsing skipped for %s", file_path.name
            )
            return {"error": "ebooklib not installed", "text": ""}
        except Exception as exc:
            logger.warning("ePub parse failed for %s: %s", file_path.name, exc)
            return {"error": str(exc), "text": ""}

    # ── HTML ──────────────────────────────────────────────────────

    def _parse_html(self, file_path: Path) -> Dict[str, Any]:
        """Parse HTML file, extract text safely via BeautifulSoup."""
        try:
            from bs4 import BeautifulSoup

            raw = file_path.read_text(encoding="utf-8", errors="ignore")
            soup = BeautifulSoup(raw, "html.parser")

            # Remove script and style elements
            for tag in soup(["script", "style", "noscript"]):
                tag.decompose()

            text = soup.get_text(separator="\n", strip=True)

            metadata: Dict[str, Any] = {}
            title_tag = soup.find("title")
            if title_tag:
                metadata["title"] = title_tag.get_text(strip=True)

            # Extract image URLs from the document
            img_urls = []
            for img in soup.find_all("img", src=True):
                img_urls.append(img["src"])
            if img_urls:
                metadata["image_urls"] = img_urls[:20]  # cap at 20

            return {
                "text": text,
                "metadata": metadata,
                "page_count": 1,
            }
        except ImportError:
            logger.warning(
                "beautifulsoup4 not installed – HTML parsing skipped for %s",
                file_path.name,
            )
            # Fallback: strip tags with regex
            raw = file_path.read_text(encoding="utf-8", errors="ignore")
            text = re.sub(r"<[^>]+>", " ", raw)
            text = re.sub(r"\s+", " ", text).strip()
            return {"text": text, "metadata": {}, "page_count": 1}

    # ── ZIP ───────────────────────────────────────────────────────

    def _parse_zip(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from text files inside a ZIP archive (max 10 files, safe)."""
        TEXT_EXTENSIONS = {
            ".txt",
            ".md",
            ".csv",
            ".json",
            ".yaml",
            ".yml",
            ".xml",
            ".log",
            ".py",
            ".js",
        }
        MAX_FILES = 10
        MAX_SINGLE_SIZE = 5 * 1024 * 1024  # 5 MB per file

        text_parts = []
        file_list = []
        metadata: Dict[str, Any] = {}

        try:
            with zipfile.ZipFile(file_path, "r") as zf:
                entries = [i for i in zf.infolist() if not i.is_dir()]
                metadata["total_entries"] = len(entries)

                extracted = 0
                for info in entries:
                    if extracted >= MAX_FILES:
                        break
                    name = info.filename
                    ext = Path(name).suffix.lower()
                    if ext not in TEXT_EXTENSIONS:
                        continue
                    if info.file_size > MAX_SINGLE_SIZE:
                        continue
                    # Security: skip paths with .. to prevent zip-slip
                    if ".." in name:
                        continue

                    try:
                        content = zf.read(name).decode("utf-8", errors="ignore")
                        if content.strip():
                            text_parts.append(f"=== {name} ===\n{content}")
                            file_list.append(name)
                            extracted += 1
                    except Exception:
                        continue

            metadata["extracted_files"] = file_list
            metadata["extracted_count"] = len(file_list)

            if not text_parts:
                return {
                    "text": f"[ZIP: {file_path.name}, no extractable text files found]",
                    "metadata": metadata,
                    "page_count": 1,
                }

            return {
                "text": "\n\n".join(text_parts),
                "metadata": metadata,
                "page_count": len(file_list),
            }
        except zipfile.BadZipFile:
            return {"error": f"Invalid ZIP file: {file_path.name}", "text": ""}

    # ── Helpers ───────────────────────────────────────────────────

    def _get_media_duration(self, file_path: Path) -> float | None:
        """Get media duration in seconds via ffprobe."""
        try:
            import subprocess

            result = subprocess.run(
                [
                    "ffprobe",
                    "-v",
                    "error",
                    "-show_entries",
                    "format=duration",
                    "-of",
                    "default=noprint_wrappers=1:nokey=1",
                    str(file_path),
                ],
                capture_output=True,
                text=True,
                timeout=10,
            )
            if result.returncode == 0 and result.stdout.strip():
                return float(result.stdout.strip())
        except (FileNotFoundError, ValueError, Exception):
            pass
        return None


# Singleton
_file_parser_service = None


def get_file_parser_service() -> FileParserService:
    global _file_parser_service
    if _file_parser_service is None:
        _file_parser_service = FileParserService()
    return _file_parser_service
