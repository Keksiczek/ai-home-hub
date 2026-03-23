"""Tests for new media features:
- FileParserService new formats (.pptx, .html, .zip, .epub)
- FileHandlerService unified metadata
- GET /api/knowledge/files listing
- DELETE /api/knowledge/files/{file_id}
- POST /api/chat/with-files
- KB retention service
"""

import io
import sys
import textwrap
import zipfile
from pathlib import Path
from typing import Any, Dict
from unittest.mock import AsyncMock, MagicMock, patch, PropertyMock

import pytest

# ── ChromaDB shim ────────────────────────────────────────────────
_chroma_mock = MagicMock()
for _mod_name in ("chromadb", "chromadb.config"):
    sys.modules.setdefault(_mod_name, _chroma_mock)

from fastapi.testclient import TestClient  # noqa: E402
from app.main import app  # noqa: E402

# ── Test fixtures ────────────────────────────────────────────────


@pytest.fixture
def client() -> TestClient:
    with (
        patch("app.services.startup_checks.run_startup_checks", new_callable=AsyncMock),
        patch(
            "app.services.vector_store_service.get_vector_store_service",
            return_value=MagicMock(),
        ),
        patch(
            "app.services.embeddings_service.get_embeddings_service",
            return_value=MagicMock(),
        ),
    ):
        yield TestClient(app, raise_server_exceptions=True)


# ── FileParserService tests ──────────────────────────────────────


class TestFileParserServiceNew:
    """Unit tests for newly added format parsers."""

    def test_parse_html_basic(self, tmp_path: Path) -> None:
        pytest.importorskip("bs4", reason="beautifulsoup4 not installed")
        from app.services.file_parser_service import FileParserService

        html = tmp_path / "test.html"
        html.write_text(
            "<html><head><title>Test Page</title></head>"
            "<body><h1>Hello</h1><p>World</p>"
            "<script>alert(1)</script></body></html>",
            encoding="utf-8",
        )
        svc = FileParserService()
        result = svc.parse_file(html)

        assert "error" not in result
        assert "Hello" in result["text"]
        assert "World" in result["text"]
        # Script content should be stripped
        assert "alert" not in result["text"]
        assert result["metadata"].get("title") == "Test Page"

    def test_parse_html_fallback_without_bs4(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """Should strip tags via regex when beautifulsoup4 is not available."""
        import builtins

        _original_import = builtins.__import__

        def _mock_import(name: str, *args: Any, **kwargs: Any) -> Any:
            if name in ("bs4", "beautifulsoup4"):
                raise ImportError("No bs4")
            return _original_import(name, *args, **kwargs)

        html = tmp_path / "test.html"
        html.write_text("<p>Hello <b>World</b></p>", encoding="utf-8")

        from app.services.file_parser_service import FileParserService

        svc = FileParserService()

        with monkeypatch.context() as m:
            m.setattr(builtins, "__import__", _mock_import)
            result = svc._parse_html(html)

        assert "Hello" in result["text"]
        assert "World" in result["text"]

    def test_parse_zip_extracts_text(self, tmp_path: Path) -> None:
        from app.services.file_parser_service import FileParserService

        zip_path = tmp_path / "archive.zip"
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.writestr("notes.txt", "Hello from zip!")
            zf.writestr("code.py", "print('Python code')")
            zf.writestr("image.png", b"\x89PNG\r\n\x1a\n")  # Not a text file

        svc = FileParserService()
        result = svc.parse_file(zip_path)

        assert "error" not in result
        assert "Hello from zip!" in result["text"]
        assert "Python code" in result["text"]
        # PNG should not be included
        assert "PNG" not in result["text"]
        assert result["metadata"]["extracted_count"] == 2

    def test_parse_zip_zip_slip_protection(self, tmp_path: Path) -> None:
        from app.services.file_parser_service import FileParserService

        zip_path = tmp_path / "evil.zip"
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.writestr("../../../etc/passwd", "root:x:0:0:root:/root:/bin/bash")
            zf.writestr("safe.txt", "safe content")

        svc = FileParserService()
        result = svc.parse_file(zip_path)

        assert "root:x:0:0" not in result.get("text", "")
        assert "safe content" in result.get("text", "")

    def test_parse_pptx_extracts_slides(self, tmp_path: Path) -> None:
        from app.services.file_parser_service import FileParserService

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            pytest.skip("python-pptx not installed")

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "Slide body text"
        prs.save(str(pptx_path))

        svc = FileParserService()
        result = svc.parse_file(pptx_path)

        assert "error" not in result
        assert "Slide Title" in result["text"]
        assert "Slide body text" in result["text"]
        assert result["page_count"] == 1

    def test_media_type_resolution(self) -> None:
        from app.services.file_parser_service import FileParserService

        assert FileParserService.MEDIA_TYPES[".mp3"] == "audio"
        assert FileParserService.MEDIA_TYPES[".mp4"] == "video"
        assert FileParserService.MEDIA_TYPES[".docx"] == "office"
        assert FileParserService.MEDIA_TYPES[".zip"] == "archive"
        assert FileParserService.MEDIA_TYPES[".txt"] == "text"


# ── FileMetadata schema test ─────────────────────────────────────


class TestFileMetadataSchema:
    def test_file_metadata_defaults(self) -> None:
        from app.models.schemas import FileMetadata

        m = FileMetadata(filename="test.pdf", filetype="application/pdf")
        assert m.media_type == "text"
        assert m.chunk_count == 0
        assert m.collection == "default"

    def test_file_metadata_audio(self) -> None:
        from app.models.schemas import FileMetadata

        m = FileMetadata(
            filename="podcast.mp3",
            filetype="audio/mpeg",
            media_type="audio",
            pages_or_duration=2700.0,
        )
        assert m.pages_or_duration == 2700.0
        assert m.media_type == "audio"


# ── FileHandlerService tests ─────────────────────────────────────


class TestFileHandlerServiceExtended:
    def test_is_supported_new_extensions(self) -> None:
        from app.services.file_handler_service import FileHandlerService

        assert FileHandlerService.is_supported("file.pptx")
        assert FileHandlerService.is_supported("audio.mp3")
        assert FileHandlerService.is_supported("video.mp4")
        assert FileHandlerService.is_supported("ebook.epub")
        assert FileHandlerService.is_supported("page.html")
        assert FileHandlerService.is_supported("archive.zip")
        assert not FileHandlerService.is_supported("unknown.xyz")

    @pytest.mark.asyncio
    async def test_process_file_html_index_mode(self, tmp_path: Path) -> None:
        from app.services.file_handler_service import FileHandlerService

        html = tmp_path / "page.html"
        html.write_text("<html><body><p>Test content for indexing.</p></body></html>")

        result = await FileHandlerService.process_file(str(html), "index")

        assert "error" not in result
        assert result["chunk_count"] > 0
        assert "Test content" in result["text"]
        assert result["file_metadata"]["media_type"] == "text"
        assert result["file_metadata"]["filename"] == "page.html"

    @pytest.mark.asyncio
    async def test_process_file_zip_analyze_mode(self, tmp_path: Path) -> None:
        from app.services.file_handler_service import FileHandlerService

        zip_path = tmp_path / "data.zip"
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.writestr("readme.txt", "This is the readme content.")

        with patch(
            "app.services.file_handler_service._generate_summary",
            new_callable=AsyncMock,
            return_value="Summary of zip.",
        ):
            result = await FileHandlerService.process_file(str(zip_path), "analyze")

        assert "error" not in result
        assert "text_preview" in result
        assert result["file_metadata"]["media_type"] == "archive"

    @pytest.mark.asyncio
    async def test_process_file_nonexistent(self) -> None:
        from app.services.file_handler_service import FileHandlerService

        result = await FileHandlerService.process_file("/nonexistent/file.txt", "index")
        assert "error" in result


# ── Knowledge files listing endpoint ─────────────────────────────


class TestKBFilesEndpoint:
    def _make_mock_vector_store(self, files_data: list) -> MagicMock:
        vs = MagicMock()
        vs.collection.get = MagicMock(
            return_value={
                "metadatas": [f["meta"] for f in files_data],
                "documents": [f.get("doc", "") for f in files_data],
            }
        )
        return vs

    def test_list_kb_files_empty(self, client: TestClient) -> None:
        vs = self._make_mock_vector_store([])
        with patch("app.routers.knowledge.get_vector_store_service", return_value=vs):
            resp = client.get("/api/knowledge/files")
        assert resp.status_code == 200
        data = resp.json()
        assert data["total"] == 0
        assert data["files"] == []

    def test_list_kb_files_groups_by_file(self, client: TestClient) -> None:
        vs = self._make_mock_vector_store(
            [
                {
                    "meta": {
                        "file_path": "/uploads/doc.pdf",
                        "file_name": "doc.pdf",
                        "collection": "default",
                        "mtime": 1700000000.0,
                        "page_count": 5,
                    }
                },
                {
                    "meta": {
                        "file_path": "/uploads/doc.pdf",
                        "file_name": "doc.pdf",
                        "collection": "default",
                        "mtime": 1700000000.0,
                        "page_count": 5,
                    }
                },
                {
                    "meta": {
                        "file_path": "/uploads/audio.mp3",
                        "file_name": "audio.mp3",
                        "collection": "default",
                        "mtime": 1700001000.0,
                        "page_count": 1,
                    }
                },
            ]
        )
        with patch("app.routers.knowledge.get_vector_store_service", return_value=vs):
            resp = client.get("/api/knowledge/files")
        assert resp.status_code == 200
        data = resp.json()
        assert data["total"] == 2
        # Find doc.pdf entry
        doc = next(f for f in data["files"] if f["file_name"] == "doc.pdf")
        assert doc["chunk_count"] == 2
        audio = next(f for f in data["files"] if f["file_name"] == "audio.mp3")
        assert audio["media_type"] == "audio"

    def test_delete_kb_file_requires_auth(self, client: TestClient) -> None:
        """Delete without API key should be rejected (403) when api_key is configured."""
        with patch("app.utils.auth.get_settings_service") as mock_svc:
            mock_svc.return_value.load.return_value = {"api_key": "secret-test-key"}
            resp = client.delete("/api/knowledge/files/some%2Ffile.pdf")
        assert resp.status_code == 403


# ── Chat with files endpoint ─────────────────────────────────────


class TestChatWithFiles:
    def test_chat_with_files_no_files(self, client: TestClient) -> None:
        """POST /chat/with-files with no files should still process message."""
        mock_llm = MagicMock()
        mock_llm.generate = AsyncMock(
            return_value=("Test reply", {"provider": "ollama"})
        )
        mock_session = MagicMock()
        mock_session.session_exists.return_value = False
        mock_session.create_session.return_value = "sess_test"
        mock_session.get_history_for_llm.return_value = []
        mock_session.save_message = MagicMock()

        with (
            patch("app.routers.chat.get_llm_service", return_value=mock_llm),
            patch("app.routers.chat.get_session_service", return_value=mock_session),
            patch(
                "app.routers.chat.enrich_message",
                new_callable=AsyncMock,
                return_value=("hello", {}),
            ),
        ):
            resp = client.post(
                "/api/chat/with-files",
                data={"message": "Hello"},
            )

        assert resp.status_code == 200
        data = resp.json()
        assert data["reply"] == "Test reply"
        assert data["meta"]["attachments"] == []

    def test_chat_with_unsupported_file(self, client: TestClient) -> None:
        """Unsupported extension should return graceful context note."""
        mock_llm = MagicMock()
        mock_llm.generate = AsyncMock(return_value=("OK", {"provider": "ollama"}))
        mock_session = MagicMock()
        mock_session.session_exists.return_value = False
        mock_session.create_session.return_value = "sess_xyz"
        mock_session.get_history_for_llm.return_value = []
        mock_session.save_message = MagicMock()

        with (
            patch("app.routers.chat.get_llm_service", return_value=mock_llm),
            patch("app.routers.chat.get_session_service", return_value=mock_session),
            patch(
                "app.routers.chat.enrich_message",
                new_callable=AsyncMock,
                return_value=("msg", {}),
            ),
        ):
            resp = client.post(
                "/api/chat/with-files",
                data={"message": "Check this"},
                files={
                    "files": (
                        "evil.xyz",
                        io.BytesIO(b"data"),
                        "application/octet-stream",
                    )
                },
            )

        assert resp.status_code == 200
        data = resp.json()
        assert "evil.xyz" in data["meta"]["attachments"]


# ── KB Retention service tests ───────────────────────────────────


class TestKBRetention:
    @pytest.mark.asyncio
    async def test_retention_run_returns_summary(self) -> None:
        from app.services.kb_retention_service import run_kb_retention

        mock_settings = MagicMock()
        mock_settings.load.return_value = {
            "knowledge_base": {"retention_days": 30, "max_size_gb": 10}
        }

        mock_vs = MagicMock()
        mock_vs.collection.get = MagicMock(return_value={"metadatas": [], "ids": []})
        mock_vs._safe_write = AsyncMock()

        with (
            patch(
                "app.services.settings_service.get_settings_service",
                return_value=mock_settings,
            ),
            patch(
                "app.services.vector_store_service.get_vector_store_service",
                return_value=mock_vs,
            ),
        ):
            result = await run_kb_retention()

        assert "deleted_old" in result
        assert "deleted_size" in result
        assert "retention_days" in result
        assert result["retention_days"] == 30

    @pytest.mark.asyncio
    async def test_retention_deletes_old_files(self) -> None:
        """Files older than retention_days should be deleted."""
        import time
        from app.services.kb_retention_service import run_kb_retention

        old_mtime = time.time() - (40 * 86400)  # 40 days ago

        mock_settings = MagicMock()
        mock_settings.load.return_value = {
            "knowledge_base": {"retention_days": 30, "max_size_gb": 10}
        }

        mock_vs = MagicMock()
        mock_vs.collection.get = MagicMock(
            return_value={
                "metadatas": [
                    {"file_path": "/kb/old_file.pdf", "mtime": old_mtime},
                    {"file_path": "/kb/old_file.pdf", "mtime": old_mtime},
                ],
                "ids": ["chunk_1", "chunk_2"],
            }
        )
        mock_vs._safe_write = AsyncMock()

        with (
            patch(
                "app.services.settings_service.get_settings_service",
                return_value=mock_settings,
            ),
            patch(
                "app.services.vector_store_service.get_vector_store_service",
                return_value=mock_vs,
            ),
        ):
            result = await run_kb_retention()

        assert result["deleted_old"] == 1
        mock_vs._safe_write.assert_awaited_once()

    def test_retention_config_endpoint(self, client: TestClient) -> None:
        mock_settings = MagicMock()
        mock_settings.load.return_value = {
            "knowledge_base": {"retention_days": 14, "max_size_gb": 5}
        }
        with patch(
            "app.routers.knowledge.get_settings_service", return_value=mock_settings
        ):
            resp = client.get("/api/knowledge/retention/config")
        assert resp.status_code == 200
        data = resp.json()
        assert data["retention_days"] == 14
        assert data["max_size_gb"] == 5
